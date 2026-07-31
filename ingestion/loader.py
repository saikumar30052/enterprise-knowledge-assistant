import re
from pathlib import Path
from typing import Any, Dict, List

import pandas as pd
from docx import Document as DocxDocument
from langchain_core.documents import Document
from pypdf import PdfReader
from pptx import Presentation


class DocumentLoader:

    def __init__(self, documents_path: str):
        self.documents_path = Path(documents_path)

    def load_documents(self) -> List[Document]:
        documents: List[Document] = []

        if not self.documents_path.exists():
            return documents

        for file in sorted(self.documents_path.iterdir()):
            suffix = file.suffix.lower()

            if suffix == ".pdf":
                documents.extend(self._load_pdf(file))
                continue

            if suffix == ".docx":
                documents.extend(self._load_docx(file))
                continue

            if suffix == ".pptx":
                documents.extend(self._load_pptx(file))
                continue

            if suffix in [".xlsx", ".xls"]:
                documents.extend(self._load_excel_rows(file)[0])

        print(f"Loaded structure-aware documents: {len(documents)}")
        print(f"Document types: {self._count_document_types(documents)}")

        return documents

    def _load_pdf(self, file: Path) -> List[Document]:
        documents: List[Document] = []
        reader = PdfReader(str(file))
        for page_number, page in enumerate(reader.pages, start=1):
            content = (page.extract_text() or "").strip()
            documents.append(
                Document(
                    page_content=content,
                    metadata={
                        "source": file.name,
                        "document_type": "pdf",
                        "page": page_number,
                        "title": file.stem,
                    },
                )
            )
        return documents

    def _load_docx(self, file: Path) -> List[Document]:
        document = DocxDocument(str(file))
        sections: List[tuple[str, List[str]]] = []
        current_heading = ""
        current_paragraphs: List[str] = []
        has_heading = False

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = str(getattr(paragraph.style, "name", "")).lower()
            is_heading = style_name.startswith("heading")
            if is_heading:
                has_heading = True
                if current_heading or current_paragraphs:
                    sections.append((current_heading, current_paragraphs))
                current_heading = text
                current_paragraphs = []
            else:
                current_paragraphs.append(text)

        if current_heading or current_paragraphs:
            sections.append((current_heading, current_paragraphs))

        documents: List[Document] = []
        if has_heading:
            for heading, paragraphs in sections:
                content_parts = ([heading] if heading else []) + paragraphs
                content = "\n\n".join(content_parts).strip()
                if content:
                    documents.append(
                        Document(
                            page_content=content,
                            metadata={
                                "source": file.name,
                                "document_type": "docx",
                                "section": heading or "Introduction",
                                "title": file.stem,
                            },
                        )
                    )
        else:
            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text:
                    documents.append(
                        Document(
                            page_content=text,
                            metadata={
                                "source": file.name,
                                "document_type": "docx",
                                "section": "",
                                "title": file.stem,
                            },
                        )
                    )
        return documents

    def _load_pptx(self, file: Path) -> List[Document]:
        presentation = Presentation(str(file))
        documents: List[Document] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            title = ""
            try:
                if slide.shapes.title is not None:
                    title = slide.shapes.title.text.strip()
            except (AttributeError, ValueError):
                pass
            content_parts: List[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_table", False):
                    table = getattr(shape, "table", None)
                    for row in getattr(table, "rows", []):
                        row_text = " | ".join(cell.text.strip() for cell in row.cells).strip()
                        if row_text:
                            content_parts.append(row_text)
                    continue
                if getattr(shape, "has_text_frame", False):
                    text = str(getattr(shape, "text", "")).strip()
                    if not text or shape == getattr(slide.shapes, "title", None):
                        continue
                    if not title and getattr(shape, "is_placeholder", False):
                        try:
                            if shape.placeholder_format.type == 1:
                                title = text
                                continue
                        except (AttributeError, ValueError):
                            pass
                    content_parts.append(text)

            notes = ""
            try:
                notes_frame = getattr(getattr(slide, "notes_slide", None), "notes_text_frame", None)
                notes_text = str(getattr(notes_frame, "text", "")).strip()
                if notes_text:
                    notes = f"Notes:\n{notes_text}"
            except (AttributeError, ValueError):
                pass

            content = [f"Slide {slide_number}", "", "Title:", title or "Untitled", "", "Content:"]
            content.extend(content_parts)
            if notes:
                content.extend(["", notes])
            documents.append(
                Document(
                    page_content="\n".join(content).strip(),
                    metadata={
                        "source": file.name,
                        "document_type": "pptx",
                        "slide": slide_number,
                        "slide_title": title or "Untitled",
                    },
                )
            )
        return documents

    def _load_excel_rows(self, file: Path) -> tuple[List[Document], Dict[str, int]]:
        """Load Excel workbooks one row at a time, creating one document per populated row."""
        excel_docs: List[Document] = []
        loaded_rows = 0
        documents_created = 0
        rows_skipped = 0

        try:
            excel_file = pd.ExcelFile(file)
        except Exception as exc:  # pragma: no cover - defensive fallback
            print(f"Unable to read Excel file {file.name}: {exc}")
            return [], {"loaded_rows": 0, "documents_created": 0, "rows_skipped": 0}

        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file, sheet_name=sheet_name, header=None)
            if df.empty:
                continue

            header_row_index = self._detect_header_row(df)
            if header_row_index is None:
                continue

            header_row = df.iloc[header_row_index].tolist()
            header_map: dict[int, str] = {}
            for index, header_name in enumerate(header_row):
                if self._is_blank_value(header_name):
                    continue
                header_map[index] = str(header_name).strip()

            for row_number, row in df.iloc[header_row_index + 1 :].iterrows():
                row_values: Dict[str, Any] = {}
                for column_index, header_name in header_map.items():
                    row_values[header_name] = self._normalize_value(row.iloc[column_index])

                if self._is_empty_row(row_values):
                    rows_skipped += 1
                    continue

                loaded_rows += 1
                document_text = self._format_row_text(row_values)
                metadata = {
                    "document_type": "xlsx",
                    "table": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Target Table Name",
                            "Target Table",
                            "Table",
                            "TARGET",
                        )
                    ) or self._normalize_metadata_value(sheet_name),
                    "target_column": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Target Column Name",
                            "Target Column",
                            "Column",
                            "TARGET",
                        )
                    ),
                    "source_table": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Source Table Name",
                            "Source Table",
                            "Source",
                            "SOURCE TABLE",
                        )
                    ),
                    "source_column": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Source Column Name",
                            "Source Column",
                            "Source",
                            "SOURCE",
                        )
                    ),
                    "business_name": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Business Name",
                            "Business Name/ Logical Name",
                            "Logical Name",
                            "Business",
                        )
                    ),
                    "business_description": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Business Description",
                            "Business Desc",
                            "Business Name/ Logical Name",
                            "Business",
                        )
                    ),
                    "datatype": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Data Type",
                            "Datatype",
                            "Type",
                        )
                    ),
                    "transformation": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Transformation",
                            "Transformation Logic",
                            "Logic",
                        )
                    ),
                    "requirement": self._normalize_metadata_value(
                        self._get_first_value(
                            row_values,
                            "Requirement",
                            "Requirement Name",
                            "Requirement Description",
                        )
                    ),
                    "row_number": int(str(row_number)) + 1,
                    "sheet": str(sheet_name),
                    "source": str(file.name),
                    "title": file.stem,
                }
                excel_docs.append(Document(page_content=document_text, metadata=metadata))
                documents_created += 1

        print(f"Total rows read: {loaded_rows}")
        print(f"Total documents created: {documents_created}")
        print(f"Rows skipped: {rows_skipped}")

        if excel_docs:
            sample_document = excel_docs[0]
            print("Example document:")
            print(sample_document.page_content)
            print("Metadata:")
            print(sample_document.metadata)

        return excel_docs, {
            "loaded_rows": loaded_rows,
            "documents_created": documents_created,
            "rows_skipped": rows_skipped,
        }

    def _count_document_types(self, documents: List[Document]) -> Dict[str, int]:
        counts: Dict[str, int] = {}
        for document in documents:
            document_type = str(document.metadata.get("document_type", "unknown"))
            counts[document_type] = counts.get(document_type, 0) + 1
        return counts

    def _normalize_value(self, value: Any) -> Any:
        if pd.isna(value):
            return ""
        return value

    def _get_first_value(self, row_values: Dict[str, Any], *column_names: str) -> Any:
        normalized_values: Dict[str, Any] = {}
        for key, value in row_values.items():
            normalized_key = self._normalize_header_name(key)
            normalized_values[normalized_key] = value

        for column_name in column_names:
            if column_name in row_values:
                value = row_values[column_name]
                if self._is_blank_value(value):
                    continue
                return value

            normalized_name = self._normalize_header_name(column_name)
            if normalized_name in normalized_values:
                value = normalized_values[normalized_name]
                if self._is_blank_value(value):
                    continue
                return value

            for normalized_key, value in normalized_values.items():
                if normalized_name in normalized_key or normalized_key in normalized_name:
                    if self._is_blank_value(value):
                        continue
                    return value

        return ""

    def _normalize_metadata_value(self, value: Any) -> str:
        if value is None:
            return ""
        if pd.isna(value):
            return ""
        return str(value).strip().lower()

    def _is_blank_value(self, value: Any) -> bool:
        if value is None:
            return True
        if pd.isna(value):
            return True
        return str(value).strip() == ""

    def _normalize_header_name(self, name: str) -> str:
        normalized = str(name).strip().lower()
        normalized = re.sub(r"[\s_\-]+", " ", normalized)
        normalized = re.sub(r"[^a-z0-9 ]", "", normalized)
        return normalized

    def _detect_header_row(self, df: pd.DataFrame) -> int | None:
        header_keywords = [
            "requirement",
            "target table name",
            "business name",
            "business desc",
            "business description",
            "target column name",
            "source table name",
            "source column name",
            "transformation",
        ]

        for row_index, row in df.iterrows():
            normalized_values = []
            for value in row.tolist():
                if self._is_blank_value(value):
                    continue
                normalized_values.append(str(value).strip().lower())

            if not normalized_values:
                continue

            matched_keywords = sum(
                1 for value in normalized_values if any(keyword in value for keyword in header_keywords)
            )
            if matched_keywords >= 3:
                return int(str(row_index))

        return None

    def _is_empty_row(self, row_values: Dict[str, Any]) -> bool:
        return all(str(value).strip() == "" for value in row_values.values())

    def _format_row_text(self, row_values: Dict[str, Any]) -> str:
        lines: List[str] = ["Excel row", ""]
        for label, value in row_values.items():
            if value is None:
                value = ""
            if pd.isna(value):
                value = ""
            lines.append(f"{label}:")
            lines.append(str(value))
            lines.append("")

        lines.append("----------------------------------")
        return "\n".join(lines).rstrip()
